
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some colors and styles
    title_color = RGBColor(0, 51, 102) # Dark Blue
    content_color = RGBColor(0, 0, 0) # Black

    # --- Slide 1: Title Page ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "ZameenLink\nML-Based Property Price Predictor"
    subtitle.text = "Capstone Project Phase-2 Review 1\n\nTeam Members:\n[Member Name 1]\n[Member Name 2]\n[Member Name 3]\n\nSupervisor:\n[Supervisor Name]"

    # --- Slide 2: Objective ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Objective"
    content = slide.placeholders[1]
    
    msg = (
        "The aim of the project is to implement a complete Machine Learning-based property price prediction system for Bhopal with integrated scam detection.\n\n"
        "Motivation:\n"
        "- Address the issue of property scams and inaccurate valuations.\n"
        "- Provide a transparent platform for buyers and sellers.\n\n"
        "Purpose:\n"
        "- Train and compare three ML models (Linear Regression, Random Forest, XGBoost).\n"
        "- Provide accurate property valuations.\n"
        "- Detect overpriced listings (Scam Detection)."
    )
    content.text = msg

    # --- Slide 3: Clarity and Significance of Problem ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Clarity and Significance"
    content = slide.placeholders[1]
    
    msg = (
        "Problem Statement:\n"
        "Real estate markets often suffer from information asymmetry, leading to scams and inflated prices. Buyers lack tools to verify fair market prices.\n\n"
        "Significance:\n"
        "- Real-time Predictions: Instant fair value estimation.\n"
        "- Scam Detection: Flags properties overpriced by >20%.\n"
        "- Data-Driven Decisions: Empowers users with ML insights."
    )
    content.text = msg

    # --- Slide 4: Feasibility and Scope ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Feasibility and Scope"
    content = slide.placeholders[1]
    
    msg = (
        "Feasibility:\n"
        "- Technical: Uses proven ML stack (Python, Scikit-learn, React).\n"
        "- Data: Synthetic data generation module ensures data availability for training.\n"
        "- Integration: Flask backend + React frontend is a standard, robust architecture.\n\n"
        "Scope:\n"
        "- Geography: Bhopal, India.\n"
        "- Features: Price Prediction, Heatmap, Scam Alerts, Comparable Listings.\n"
        "- Platforms: Web Application (Responsive)."
    )
    content.text = msg

    # --- Slide 5: Existing Projects ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Existing Projects"
    content = slide.placeholders[1]
    
    msg = (
        "Existing Solutions:\n"
        "- MagicBricks, 99Acres (Listing platforms).\n\n"
        "Limitations of Existing Systems:\n"
        "- Focus on listings, not fair price validation.\n"
        "- Lack of automated scam detection based on price deviation.\n\n"
        "ZameenLink Differences:\n"
        "- ML-based Fair Price Prediction.\n"
        "- Active Scam Detection (Critical, High, Medium Risk flags).\n"
        "- Detailed Model Performance Metrics for transparency."
    )
    content.text = msg

    # --- Slide 6: Project Modules ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Project Modules"
    content = slide.placeholders[1]

    msg = (
        "1. Backend Processing:\n"
        "   - Flask API, Data Preprocessing, ML Model Training (RF, XGB, LR).\n"
        "   - Scam Detection Logic.\n\n"
        "2. Frontend Interface:\n"
        "   - React.js UI, Leaflet Interactive Map, Recharts for Analytics.\n"
        "   - Search & Prediction Panel.\n\n"
        "3. Data Engine:\n"
        "   - Synthetic Data Generator (Bhopal specific).\n"
        "   - Feature Engineering (Distance to landmarks, etc.)."
    )
    content.text = msg
    
    # --- Slide 7: References ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "References"
    content = slide.placeholders[1]

    msg = (
        "Technologies Used:\n"
        "- Scikit-learn, XGBoost (Machine Learning)\n"
        "- Flask, Python (Backend)\n"
        "- React, Leaflet, Tailwind CSS (Frontend)\n"
        "- Pandas, NumPy (Data Processing)\n\n"
        "Methodologies:\n"
        "- Random Forest Regression\n"
        "- Gradient Boosting (XGBoost)\n"
        "- Geospatial Analysis"
    )
    content.text = msg

    # Save
    prs.save('d:/Projects/zameenlink/Phase1_Report/ZameenLink_Presentation.pptx')
    print("Presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
